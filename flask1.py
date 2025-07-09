import os
import sys
import io
import cv2
import numpy as np
import streamlit as st
from PIL import Image
from rembg import remove
from pptx import Presentation
from pptx.util import Inches
from ultralytics import YOLO
import torch
# -------------------------------------------------------------------
# Helper: resource_path()
# Returns the absolute path to a resource whether running as a script
# or as a bundled executable via PyInstaller.
def resource_path(relative_path):
    if getattr(sys, 'frozen', False):  # running in a bundle
        base_path = sys._MEIPASS
    else:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)


# -------------------------------------------------------------------
# Use resource_path to get the correct paths for the resource files

model_resource = resource_path("best.pt")

# Load the PPTX template into memory using the bundled resource.


# Global variables that the processing functions rely on.
patient_name = ""
folder_path = ""


# -------------------------------------------------------------------
# Processing Functions

def create_new_presentation():

    pptx_resource = "main.pptx" # the presentation template
    with open(pptx_resource, "rb") as file:
        pptx_data = file.read()
    pptx_stream = io.BytesIO(pptx_data)
    new_presentation = Presentation(pptx_stream)


def copy_and_rename_convert_images():
    """
    Copies and renames only the images in the uploaded_image_paths dictionary,
    converting them to JPG.
    """
    for var_name, image_path in uploaded_image_paths.items():
        if os.path.isfile(image_path) and image_path.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.tiff')):
            try:
                copied_image_name = f"{var_name}.jpg"
                copied_image_path = os.path.join(folder_path, copied_image_name)
                image = Image.open(image_path).convert("RGB")
                image.save(copied_image_path, "JPEG")
                st.write(f"Copied, renamed, and converted: {copied_image_name}")
            except Exception as e:
                st.write(f"Error processing {image_path}: {e}")



def remove_background(image_name):
    """Removes the background from an image and overwrites the file."""
    image_path = os.path.join(folder_path, image_name)
    if os.path.isfile(image_path) and image_name.lower().endswith(('.png', '.jpg', '.jpeg')):
        try:
            input_image = Image.open(image_path)
            output_image = remove(input_image)
            white_background = Image.new("RGB", output_image.size, (255, 255, 255))
            white_background.paste(output_image, mask=output_image.split()[3])
            white_background.save(image_path)
            st.write(f"Background removed: {image_name}")
        except Exception as e:
            st.write(f"Error processing {image_name}: {e}")


def crop_personal(image_name, left_padding=0.2, right_padding=0.2, above_padding=0.3, bottom_padding=0.1):
    """Crops personal images using face detection."""
    # Adjust the face detector paths as needed. Here we use absolute paths from your desktop version.
    face_detector_pb = r'D:\AI cours by Hasoob\Real_Projects\Ortho_Presentation_Creater_desktop\.venv\Lib\site-packages\cv2\data\opencv_face_detector_uint8.pb'
    face_detector_pbtxt = r'D:\AI cours by Hasoob\Real_Projects\Ortho_Presentation_Creater_desktop\.venv\Lib\site-packages\cv2\data\opencv_face_detector.pbtxt'

    net = cv2.dnn.readNetFromTensorflow(face_detector_pb, face_detector_pbtxt)
    image_path = os.path.join(folder_path, image_name)
    img = cv2.imread(image_path)
    if img is None:
        st.write(f"Error: Unable to read image {image_path}")
        return image_path
    h, w = img.shape[:2]
    blob = cv2.dnn.blobFromImage(img, scalefactor=1.0, size=(300, 300), mean=(104.0, 177.0, 123.0))
    net.setInput(blob)
    detections = net.forward()
    best_face = None
    for i in range(detections.shape[2]):
        confidence = detections[0, 0, i, 2]
        if confidence > 0.5:
            box = detections[0, 0, i, 3:7] * np.array([w, h, w, h])
            x, y, x_max, y_max = box.astype("int")
            best_face = (x, y, x_max - x, y_max - y)
    if best_face:
        x, y, box_w, box_h = best_face
        top_pad = int(box_h * above_padding)
        bottom_pad = int(box_h * bottom_padding)
        left_pad = int(box_w * left_padding)
        right_pad = int(box_w * right_padding)
        y = max(0, y - top_pad)
        box_h = min(img.shape[0] - y, box_h + top_pad + bottom_pad)
        x = max(0, x - left_pad)
        box_w = min(img.shape[1] - x, box_w + left_pad + right_pad)
        cropped_img = img[y:y + box_h, x:x + box_w]
        cv2.imwrite(image_path, cropped_img)
        st.write(f"Photo cropped: {image_name}")
        return image_path
    else:
        st.write("No face detected")
        return image_path


def crop_arch(image_name):
    """Crops arch images using a YOLO model."""
    image_path = os.path.join(folder_path, image_name)


    # Load YOLO model
    #model = YOLO(model_resource)
    model = YOLO(model_resource)

    image = cv2.imread(image_path)
    if image is None:
        st.write(f"Error: Unable to read image {image_name}")
        return
    results = model(image_path, verbose=False)
    if results and results[0].boxes.xyxy.shape[0] > 0:
        x1, y1, x2, y2 = map(int, results[0].boxes.xyxy[0])
        cropped_object = image[y1:y2, x1:x2]
        cv2.imwrite(image_path, cropped_object)
        st.write(f"Photo cropped: {image_name}")
    else:
        st.write(f"No objects detected in {image_name}")


def resize_image(image_name, max_width=None, max_height=None, file_prefix=''):
    """Resizes an image to given dimensions and saves a new version with an optional prefix."""
    image_path = os.path.join(folder_path, image_name)
    img = Image.open(image_path)
    orig_width, orig_height = img.size
    dpi = 96
    max_width_px = max_width * dpi if max_width else None
    max_height_px = max_height * dpi if max_height else None
    if max_width_px and max_height_px:
        scale = min(max_width_px / orig_width, max_height_px / orig_height)
    elif max_width_px:
        scale = max_width_px / orig_width
    elif max_height_px:
        scale = max_height_px / orig_height
    else:
        scale = 1
    new_width = int(orig_width * scale)
    new_height = int(orig_height * scale)
    resized_img = img.resize((new_width, new_height), Image.LANCZOS)
    base_name, ext = os.path.splitext(image_name)
    new_file_name = f"{file_prefix}{base_name}{ext}"
    temp_path = os.path.join(folder_path, new_file_name)
    resized_img.save(temp_path)
    st.write(f"Resized image: {new_file_name}")
    return new_file_name


def insert_image(slide_index, image_name, left=None, bottom=None, right=None, top=None):
    """Inserts an image into a slide of the PowerPoint presentation."""
    pptx_path_local = os.path.join(folder_path, os.path.splitext(patient_name.replace(' ', '_'))[0] + '.pptx')
    ppt = Presentation(pptx_path_local)
    slide = ppt.slides[slide_index]
    image_path = os.path.join(folder_path, image_name)
    with Image.open(image_path) as img:
        dpi_info = img.info.get('dpi', (96, 96))
        # Robust handling of dpi_info: if it's not a proper 2-tuple, default to (96,96)
        if isinstance(dpi_info, tuple):
            if len(dpi_info) < 2:
                if len(dpi_info) == 1:
                    dpi = (dpi_info[0], dpi_info[0])
                else:
                    dpi = (96, 96)
            else:
                dpi = dpi_info
        elif isinstance(dpi_info, (int, float)):
            dpi = (dpi_info, dpi_info)
        else:
            dpi = (96, 96)
        width_in_inches = img.width / dpi[0]
        height_in_inches = img.height / dpi[1]
    slide_width = ppt.slide_width
    slide_height = ppt.slide_height
    if right is not None:
        left_val = slide_width - Inches(right) - Inches(width_in_inches)
    elif left is not None:
        left_val = Inches(left)
    else:
        left_val = Inches(0)
    if bottom is not None:
        top_val = slide_height - Inches(bottom) - Inches(height_in_inches)
    elif top is not None:
        top_val = Inches(top)
    else:
        top_val = Inches(0)
    slide.shapes.add_picture(image_path, left_val, top_val,
                             width=Inches(width_in_inches),
                             height=Inches(height_in_inches))
    ppt.save(pptx_path_local)
    st.write(f"Inserted image: {image_name}")


def run_all_processing():
    """Runs the complete processing pipeline."""
    # Step 1: Copy, rename, and convert images
    copy_and_rename_convert_images()

    # Step 2: Remove background from personal images
    for img in ['pre_personal_front.jpg', 'pre_personal_smile.jpg', 'pre_personal_oblique.jpg',
                'pre_personal_profile.jpg', 'post_personal_front.jpg', 'post_personal_smile.jpg',
                'post_personal_oblique.jpg', 'post_personal_profile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            remove_background(img)

    # Step 3: Crop personal images
    for img in ['pre_personal_front.jpg', 'pre_personal_smile.jpg', 'post_personal_front.jpg',
                'post_personal_smile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_personal(img)
    for img in ['pre_personal_oblique.jpg', 'post_personal_oblique.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_personal(img, left_padding=0.5)
    for img in ['pre_personal_profile.jpg', 'post_personal_profile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_personal(img, left_padding=1)

    # Step 4: Crop arch images
    for img in ['pre_arch_right.jpg', 'pre_arch_front.jpg', 'pre_arch_left.jpg', 'pre_arch_upper.jpg',
                'pre_arch_lower.jpg', 'post_arch_right.jpg', 'post_arch_front.jpg', 'post_arch_left.jpg',
                'post_arch_upper.jpg', 'post_arch_lower.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            crop_arch(img)

    # Step 5: Resize images for personal and arch images
    for img in ['pre_personal_front.jpg', 'pre_personal_smile.jpg', 'pre_personal_oblique.jpg',
                'pre_personal_profile.jpg', 'post_personal_front.jpg', 'post_personal_smile.jpg',
                'post_personal_oblique.jpg', 'post_personal_profile.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            resize_image(img, 2.9, 3.9)
            resize_image(img, 2.31, 3.01, 'g_')
    for img in ['pre_arch_right.jpg', 'pre_arch_front.jpg', 'pre_arch_left.jpg',
                'post_arch_right.jpg', 'post_arch_front.jpg', 'post_arch_left.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            resize_image(img, 3.53, 2.17)
            resize_image(img, 2.7, 1.77, 'g_')
    for img in ['pre_arch_upper.jpg', 'pre_arch_lower.jpg', 'post_arch_upper.jpg', 'post_arch_lower.jpg']:
        if os.path.exists(os.path.join(folder_path, img)):
            resize_image(img, 3.74, 2.76)
            resize_image(img, 2.69, 1.95, 'g_')

    # New Processing for Cast and X-ray Images
    for key in ["pre_cast_right", "pre_cast_front", "pre_cast_left"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.53, 2.17)
    for key in ["pre_cast_upper", "pre_cast_lower"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.74, 2.76)
    for key in ["post_cast_right", "post_cast_front", "post_cast_left"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.53, 2.17)
    for key in ["post_cast_upper", "post_cast_lower"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 3.74, 2.76)
    for key in ["pre-Panoramic X-Ray"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 9.5, 4.85)
    for key in ["pre-Cephalometric X-Ray", "pre-Cephalometric X-Ray Tracing"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 5.97, 4.96)
    for key in ["post-Panoramic X-Ray"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 9.5, 4.85)
    for key in ["post-Cephalometric X-Ray", "post-Cephalometric X-Ray Tracing"]:
        filename = f"{key}.jpg"
        if os.path.exists(os.path.join(folder_path, filename)):
            resize_image(filename, 5.97, 4.96)

    # Step 6: Insert images into PowerPoint

    # Personal Images (Pre: slide index 2, Post: slide index 9)
    if os.path.exists(os.path.join(folder_path, 'pre_personal_front.jpg')):
        insert_image(2, 'pre_personal_front.jpg', 0.59, 1.44)
    if os.path.exists(os.path.join(folder_path, 'pre_personal_smile.jpg')):
        insert_image(2, 'pre_personal_smile.jpg', 3.67, 1.44)
    if os.path.exists(os.path.join(folder_path, 'pre_personal_oblique.jpg')):
        insert_image(2, 'pre_personal_oblique.jpg', 6.72, 1.44)
    if os.path.exists(os.path.join(folder_path, 'pre_personal_profile.jpg')):
        insert_image(2, 'pre_personal_profile.jpg', 9.78, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_front.jpg')):
        insert_image(9, 'post_personal_front.jpg', 0.6, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_smile.jpg')):
        insert_image(9, 'post_personal_smile.jpg', 3.67, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_oblique.jpg')):
        insert_image(9, 'post_personal_oblique.jpg', 6.72, 1.44)
    if os.path.exists(os.path.join(folder_path, 'post_personal_profile.jpg')):
        insert_image(9, 'post_personal_profile.jpg', 9.79, 1.44)

    # Arch Images (Pre: slide index 3, Post: slide index 10)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_right.jpg')):
        insert_image(3, 'pre_arch_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_front.jpg')):
        insert_image(3, 'pre_arch_front.jpg', 4.87, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_left.jpg')):
        insert_image(3, 'pre_arch_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_upper.jpg')):
        insert_image(3, 'pre_arch_upper.jpg', 2.8, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'pre_arch_lower.jpg')):
        insert_image(3, 'pre_arch_lower.jpg', 6.74, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'post_arch_right.jpg')):
        insert_image(10, 'post_arch_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_arch_front.jpg')):
        insert_image(10, 'post_arch_front.jpg', 4.86, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_arch_left.jpg')):
        insert_image(10, 'post_arch_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_arch_upper.jpg')):
        insert_image(10, 'post_arch_upper.jpg', 2.79, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'post_arch_lower.jpg')):
        insert_image(10, 'post_arch_lower.jpg', 6.74, top=4.05)

    # Personal and Arch Images (grouped)(Pre: slide index 4, Post: slide index 11)
    # pre personal
    if os.path.exists(os.path.join(folder_path, 'g_pre_personal_front.jpg')):
        insert_image(4, 'g_pre_personal_front.jpg', 0.9, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_pre_personal_smile.jpg')):
        insert_image(4, 'g_pre_personal_smile.jpg', 3.98, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_pre_personal_oblique.jpg')):
        insert_image(4, 'g_pre_personal_oblique.jpg', 7.03, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_pre_personal_profile.jpg')):
        insert_image(4, 'g_pre_personal_profile.jpg', 10.09, 4.2)

    # post personal
    if os.path.exists(os.path.join(folder_path, 'g_post_personal_front.jpg')):
        insert_image(11, 'g_post_personal_front.jpg', 0.9, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_post_personal_smile.jpg')):
        insert_image(11, 'g_post_personal_smile.jpg', 3.98, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_post_personal_oblique.jpg')):
        insert_image(11, 'g_post_personal_oblique.jpg', 7.03, 4.2)
    if os.path.exists(os.path.join(folder_path, 'g_post_personal_profile.jpg')):
        insert_image(11, 'g_post_personal_profile.jpg', 10.09, 4.2)

    # pre arch
    if os.path.exists(os.path.join(folder_path, 'g_pre_arch_right.jpg')):
        insert_image(4, 'g_pre_arch_right.jpg', 2.01, top=3.46)
    if os.path.exists(os.path.join(folder_path, 'g_pre_arch_front.jpg')):
        insert_image(4, 'g_pre_arch_front.jpg', 4.93, top=3.45)
    if os.path.exists(os.path.join(folder_path, 'g_pre_arch_left.jpg')):
        insert_image(4, 'g_pre_arch_left.jpg', 7.84, top=3.45)
    if os.path.exists(os.path.join(folder_path, 'g_pre_arch_upper.jpg')):
        insert_image(4, 'g_pre_arch_upper.jpg', 3.42, top=5.35)
    if os.path.exists(os.path.join(folder_path, 'g_pre_arch_lower.jpg')):
        insert_image(4, 'g_pre_arch_lower.jpg', 6.44, top=5.35)

    # post arch
    if os.path.exists(os.path.join(folder_path, 'g_post_arch_right.jpg')):
        insert_image(11, 'g_post_arch_right.jpg', 2.01, top=3.46)
    if os.path.exists(os.path.join(folder_path, 'g_post_arch_front.jpg')):
        insert_image(11, 'g_post_arch_front.jpg', 4.93, top=3.45)
    if os.path.exists(os.path.join(folder_path, 'g_post_arch_left.jpg')):
        insert_image(11, 'g_post_arch_left.jpg', 7.84, top=3.45)
    if os.path.exists(os.path.join(folder_path, 'g_post_arch_upper.jpg')):
        insert_image(11, 'g_post_arch_upper.jpg', 3.42, top=5.35)
    if os.path.exists(os.path.join(folder_path, 'g_post_arch_lower.jpg')):
        insert_image(11, 'g_post_arch_lower.jpg', 6.44, top=5.35)

    # Pre-Cast Images (insert into slide index 5)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_right.jpg')):
        insert_image(5, 'pre_cast_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_front.jpg')):
        insert_image(5, 'pre_cast_front.jpg', 4.87, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_left.jpg')):
        insert_image(5, 'pre_cast_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_upper.jpg')):
        insert_image(5, 'pre_cast_upper.jpg', 2.8, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'pre_cast_lower.jpg')):
        insert_image(5, 'pre_cast_lower.jpg', 6.74, top=4.05)

    # Post-Cast Images (insert into slide index 12)
    if os.path.exists(os.path.join(folder_path, 'post_cast_right.jpg')):
        insert_image(12, 'post_cast_right.jpg', 1.09, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_cast_front.jpg')):
        insert_image(12, 'post_cast_front.jpg', 4.87, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_cast_left.jpg')):
        insert_image(12, 'post_cast_left.jpg', 8.63, top=1.7)
    if os.path.exists(os.path.join(folder_path, 'post_cast_upper.jpg')):
        insert_image(12, 'post_cast_upper.jpg', 2.8, top=4.05)
    if os.path.exists(os.path.join(folder_path, 'post_cast_lower.jpg')):
        insert_image(12, 'post_cast_lower.jpg', 6.74, top=4.05)

    # Pre-X-rays Images (only resized; insertion based on type)
    if os.path.exists(os.path.join(folder_path, 'pre-Panoramic X-Ray.jpg')):
        insert_image(6, 'pre-Panoramic X-Ray.jpg', left=1.92, top=1.33)
    if os.path.exists(os.path.join(folder_path, 'pre-Cephalometric X-Ray.jpg')):
        insert_image(7, 'pre-Cephalometric X-Ray.jpg', left=0.55, top=1.5)
    if os.path.exists(os.path.join(folder_path, 'pre-Cephalometric X-Ray Tracing.jpg')):
        insert_image(7, 'pre-Cephalometric X-Ray Tracing.jpg', left=6.81, top=1.5)

    # Post-X-rays Images
    if os.path.exists(os.path.join(folder_path, 'post-Panoramic X-Ray.jpg')):
        insert_image(13, 'post-Panoramic X-Ray.jpg', left=1.92, top=1.33)
    if os.path.exists(os.path.join(folder_path, 'post-Cephalometric X-Ray.jpg')):
        insert_image(14, 'post-Cephalometric X-Ray.jpg', left=0.55, top=1.5)
    if os.path.exists(os.path.join(folder_path, 'post-Cephalometric X-Ray Tracing.jpg')):
        insert_image(14, 'post-Cephalometric X-Ray Tracing.jpg', left=6.81, top=1.5)



