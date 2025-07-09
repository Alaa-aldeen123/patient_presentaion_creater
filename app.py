import os
import sys
import io
import cv2
import numpy as np
import zipfile
from flask import Flask, render_template, request, send_file, redirect, url_for, session
from werkzeug.utils import secure_filename
from PIL import Image
from rembg import remove
from pptx import Presentation
from pptx.util import Inches
from ultralytics import YOLO
import torch

app = Flask(__name__)
app.secret_key = '12345'


# Helper function for resource paths
def resource_path(relative_path):
    if getattr(sys, 'frozen', False):
        base_path = sys._MEIPASS
    else:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)


# Load models and templates
pptx_resource = resource_path("main.pptx")
model_resource = resource_path("best.pt")
face_detector_pb = resource_path(".venv/Lib/site-packages/cv2/data/opencv_face_detector_uint8.pb")
face_detector_pbtxt = resource_path(".venv/Lib/site-packages/cv2/data/opencv_face_detector.pbtxt")

# Load resources once at startup
with open(pptx_resource, "rb") as file:
    pptx_template = file.read()

# Load models
face_net = cv2.dnn.readNetFromTensorflow(face_detector_pb, face_detector_pbtxt)
yolo_model = YOLO(model_resource)

# Define categories
CATEGORIES = {
    "Pre-Personal Images": ["pre_personal_front", "pre_personal_smile", "pre_personal_oblique", "pre_personal_profile"],
    "Pre-Arch Images": ["pre_arch_right", "pre_arch_front", "pre_arch_left", "pre_arch_upper", "pre_arch_lower"],
    "Pre-Cast Images": ["pre_cast_right", "pre_cast_front", "pre_cast_left", "pre_cast_upper", "pre_cast_lower"],
    "Pre-X-rays Images": ["pre-Panoramic_X-Ray", "pre-Cephalometric_X-Ray", "pre-Cephalometric_X-Ray_Tracing"],
    "Post-Personal Images": ["post_personal_front", "post_personal_smile", "post_personal_oblique","post_personal_profile"],
    "Post-Arch Images": ["post_arch_right", "post_arch_front", "post_arch_left", "post_arch_upper", "post_arch_lower"],
    "Post-Cast Images": ["post_cast_right", "post_cast_front", "post_cast_left", "post_cast_upper", "post_cast_lower"],
    "Post-X-rays Images": ["post-Panoramic_X-Ray", "post-Cephalometric_X-Ray", "post-Cephalometric_X-Ray_Tracing"]
}


# Processing functions (in-memory versions)
def create_presentation():
    return Presentation(io.BytesIO(pptx_template))


def convert_image(img_bytes, format='JPEG'):
    img = Image.open(io.BytesIO(img_bytes))
    if img.mode != 'RGB':
        img = img.convert('RGB')
    output = io.BytesIO()
    img.save(output, format=format)
    return output.getvalue()


def remove_bg(img_bytes):
    img = Image.open(io.BytesIO(img_bytes))
    output_img = remove(img)
    white_bg = Image.new("RGB", output_img.size, (255, 255, 255))
    white_bg.paste(output_img, mask=output_img.split()[3] if output_img.mode == 'RGBA' else None)
    output = io.BytesIO()
    white_bg.save(output, format='JPEG')
    return output.getvalue()


def crop_personal(img_bytes, left_padding=0.2, right_padding=0.2, above_padding=0.3, bottom_padding=0.1):
    img_arr = np.frombuffer(img_bytes, np.uint8)
    img = cv2.imdecode(img_arr, cv2.IMREAD_COLOR)
    h, w = img.shape[:2]

    # Create blob and detect faces
    blob = cv2.dnn.blobFromImage(img, 1.0, (300, 300), (104.0, 177.0, 123.0))
    face_net.setInput(blob)
    detections = face_net.forward()

    best_face = None
    for i in range(detections.shape[2]):
        confidence = detections[0, 0, i, 2]
        if confidence > 0.5:
            box = detections[0, 0, i, 3:7] * np.array([w, h, w, h])
            x, y, x_max, y_max = box.astype("int")
            best_face = (x, y, x_max - x, y_max - y)
            break  # Use first detected face

    if best_face:
        x, y, box_w, box_h = best_face
        top_pad = int(box_h * above_padding)
        bottom_pad = int(box_h * bottom_padding)
        left_pad = int(box_w * left_padding)
        right_pad = int(box_w * right_padding)

        y = max(0, y - top_pad)
        box_h = min(img.shape[0] - y, box_h + top_pad + bottom_pad)
        x = max(0, x - left_pad)
        box_w = min(img.shape[1] - x, box_w + left_pad + right_pad)

        cropped_img = img[y:y + box_h, x:x + box_w]
        _, img_encoded = cv2.imencode('.jpg', cropped_img)
        return img_encoded.tobytes()
    return img_bytes


def crop_arch(img_bytes):
    img_arr = np.frombuffer(img_bytes, np.uint8)
    img = cv2.imdecode(img_arr, cv2.IMREAD_COLOR)

    # Use YOLO model
    results = yolo_model(img, verbose=False)

    if results and results[0].boxes.xyxy.shape[0] > 0:
        x1, y1, x2, y2 = map(int, results[0].boxes.xyxy[0])
        cropped_img = img[y1:y2, x1:x2]
        _, img_encoded = cv2.imencode('.jpg', cropped_img)
        return img_encoded.tobytes()
    return img_bytes


def resize_image(img_bytes, max_width=None, max_height=None):
    img = Image.open(io.BytesIO(img_bytes))
    orig_width, orig_height = img.size
    dpi = 96
    max_width_px = max_width * dpi if max_width else None
    max_height_px = max_height * dpi if max_height else None

    if max_width_px and max_height_px:
        scale = min(max_width_px / orig_width, max_height_px / orig_height)
    elif max_width_px:
        scale = max_width_px / orig_width
    elif max_height_px:
        scale = max_height_px / orig_height
    else:
        scale = 1

    new_width = int(orig_width * scale)
    new_height = int(orig_height * scale)
    resized_img = img.resize((new_width, new_height), Image.LANCZOS)

    output = io.BytesIO()
    resized_img.save(output, format='JPEG')
    return output.getvalue()


def insert_image(presentation, slide_index, img_bytes, left=None, bottom=None, right=None, top=None):
    slide = presentation.slides[slide_index]
    img_stream = io.BytesIO(img_bytes)

    # Add image to slide
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Calculate position
    with Image.open(img_stream) as img:
        dpi = img.info.get('dpi', (96, 96))[0]
        width_in_inches = img.width / dpi
        height_in_inches = img.height / dpi

        if right is not None:
            left_val = slide_width - Inches(right) - Inches(width_in_inches)
        elif left is not None:
            left_val = Inches(left)
        else:
            left_val = Inches(0)

        if bottom is not None:
            top_val = slide_height - Inches(bottom) - Inches(height_in_inches)
        elif top is not None:
            top_val = Inches(top)
        else:
            top_val = Inches(0)

    img_stream.seek(0)
    slide.shapes.add_picture(
        img_stream, left_val, top_val,
        width=Inches(width_in_inches),
        height=Inches(height_in_inches)
    )
    return presentation


def process_all_images(patient_name, uploaded_files):
    # Create in-memory storage for processed files
    processed_files = {}

    # 1. Convert and store all uploaded images
    for key, file_data in uploaded_files.items():
        if file_data:
            processed_files[f"{key}.jpg"] = convert_image(file_data)

    # 2. Remove backgrounds from personal images
    personal_images = [
        'pre_personal_front.jpg', 'pre_personal_smile.jpg', 'pre_personal_oblique.jpg', 'pre_personal_profile.jpg',
        'post_personal_front.jpg', 'post_personal_smile.jpg', 'post_personal_oblique.jpg', 'post_personal_profile.jpg'
    ]
    for img in personal_images:
        if img in processed_files:
            processed_files[img] = remove_bg(processed_files[img])

    # 3. Crop personal images
    for img in ['pre_personal_front.jpg', 'pre_personal_smile.jpg', 'post_personal_front.jpg',
                'post_personal_smile.jpg']:
        if img in processed_files:
            processed_files[img] = crop_personal(processed_files[img])

    for img in ['pre_personal_oblique.jpg', 'post_personal_oblique.jpg']:
        if img in processed_files:
            processed_files[img] = crop_personal(processed_files[img], left_padding=0.5)

    for img in ['pre_personal_profile.jpg', 'post_personal_profile.jpg']:
        if img in processed_files:
            processed_files[img] = crop_personal(processed_files[img], left_padding=1)

    # 4. Crop arch images
    arch_images = [
        'pre_arch_right.jpg', 'pre_arch_front.jpg', 'pre_arch_left.jpg', 'pre_arch_upper.jpg', 'pre_arch_lower.jpg',
        'post_arch_right.jpg', 'post_arch_front.jpg', 'post_arch_left.jpg', 'post_arch_upper.jpg', 'post_arch_lower.jpg'
    ]
    for img in arch_images:
        if img in processed_files:
            processed_files[img] = crop_arch(processed_files[img])

    # 5. Resize images
    # Personal images
    for img in personal_images:
        if img in processed_files:
            processed_files[f"resized_{img}"] = resize_image(processed_files[img], 2.9, 3.9)
            processed_files[f"g_resized_{img}"] = resize_image(processed_files[img], 2.31, 3.01)

    # Arch images
    for img in ['pre_arch_right.jpg', 'pre_arch_front.jpg', 'pre_arch_left.jpg',
                'post_arch_right.jpg', 'post_arch_front.jpg', 'post_arch_left.jpg']:
        if img in processed_files:
            processed_files[f"resized_{img}"] = resize_image(processed_files[img], 3.53, 2.17)
            processed_files[f"g_resized_{img}"] = resize_image(processed_files[img], 2.7, 1.77)

    for img in ['pre_arch_upper.jpg', 'pre_arch_lower.jpg', 'post_arch_upper.jpg', 'post_arch_lower.jpg']:
        if img in processed_files:
            processed_files[f"resized_{img}"] = resize_image(processed_files[img], 3.74, 2.76)
            processed_files[f"g_resized_{img}"] = resize_image(processed_files[img], 2.69, 1.95)

    # casts
    for img in ['pre_cast_right.jpg', 'pre_cast_front.jpg', 'pre_cast_left.jpg',
                'post_cast_right.jpg', 'post_cast_front.jpg', 'post_cast_left.jpg']:
        if img in processed_files:
            processed_files[f"resized_{img}"] = resize_image(processed_files[img], 3.53, 2.17)

    for img in ['pre_cast_upper.jpg', 'pre_cast_lower.jpg', 'post_cast_upper.jpg', 'post_cast_lower.jpg']:
        if img in processed_files:
            processed_files[f"resized_{img}"] = resize_image(processed_files[img], 3.74, 2.76)

    # x-rays
    for img in ["pre-Panoramic_X-Ray.jpg", "post-Panoramic_X-Ray.jpg"]:
        if img in processed_files:
            processed_files[f"resized_{img}"] = resize_image(processed_files[img], 9.5, 4.85)

    for img in ["pre-Cephalometric_X-Ray.jpg", "pre-Cephalometric_X-Ray_Tracing.jpg", "post-Cephalometric_X-Ray.jpg", "post-Cephalometric_X-Ray_Tracing.jpg"]:
        if img in processed_files:
            processed_files[f"resized_{img}"] = resize_image(processed_files[img], 5.9, 4.96)


    # 6. Create presentation
    presentation = create_presentation()

    # Insert images into presentation

    # Personal Images (Pre: slide index 2, Post: slide index 9)
    # pre personal
    if 'resized_pre_personal_front.jpg' in processed_files:
        presentation = insert_image(presentation, 2, processed_files['resized_pre_personal_front.jpg'], 0.59, 1.44)
    if 'resized_pre_personal_smile.jpg' in processed_files:
        presentation = insert_image(presentation, 2, processed_files['resized_pre_personal_smile.jpg'], 3.67, 1.44)
    if 'resized_pre_personal_oblique.jpg' in processed_files:
        presentation = insert_image(presentation, 2, processed_files['resized_pre_personal_oblique.jpg'], 6.72,
                                    1.44)
    if 'resized_pre_personal_profile.jpg' in processed_files:
        presentation = insert_image(presentation, 2, processed_files['resized_pre_personal_profile.jpg'], 9.78,
                                    1.44)

    # post personal
    if 'resized_post_personal_front.jpg' in processed_files:
        presentation = insert_image(presentation, 9, processed_files['resized_post_personal_front.jpg'], 0.6, 1.44)
    if 'resized_post_personal_smile.jpg' in processed_files:
        presentation = insert_image(presentation, 9, processed_files['resized_post_personal_smile.jpg'], 3.67, 1.44)
    if 'resized_post_personal_oblique.jpg' in processed_files:
        presentation = insert_image(presentation, 9, processed_files['resized_post_personal_oblique.jpg'], 6.72,
                                    1.44)
    if 'resized_post_personal_profile.jpg' in processed_files:
        presentation = insert_image(presentation, 9, processed_files['resized_post_personal_profile.jpg'], 9.79,
                                    1.44)

    # Arch Images (Pre: slide index 3, Post: slide index 10)
    # pre arches
    if 'resized_pre_arch_right.jpg' in processed_files:
        presentation = insert_image(presentation, 3, processed_files['resized_pre_arch_right.jpg'], top=1.7,
                                    left=1.09)
    if 'resized_pre_arch_front.jpg' in processed_files:
        presentation = insert_image(presentation, 3, processed_files['resized_pre_arch_front.jpg'], top=1.7,
                                    left=4.87)
    if 'resized_pre_arch_left.jpg' in processed_files:
        presentation = insert_image(presentation, 3, processed_files['resized_pre_arch_left.jpg'], top=1.7,
                                    left=8.63)
    if 'resized_pre_arch_upper.jpg' in processed_files:
        presentation = insert_image(presentation, 3, processed_files['resized_pre_arch_upper.jpg'], top=4.05,
                                    left=2.8)
    if 'resized_pre_arch_lower.jpg' in processed_files:
        presentation = insert_image(presentation, 3, processed_files['resized_pre_arch_lower.jpg'], top=4.05,
                                    left=6.74)

    # post arches
    if 'resized_post_arch_right.jpg' in processed_files:
        presentation = insert_image(presentation, 10, processed_files['resized_post_arch_right.jpg'], top=1.7,
                                    left=1.09)
    if 'resized_post_arch_front.jpg' in processed_files:
        presentation = insert_image(presentation, 10, processed_files['resized_post_arch_front.jpg'], top=1.7,
                                    left=4.86)
    if 'resized_post_arch_left.jpg' in processed_files:
        presentation = insert_image(presentation, 10, processed_files['resized_post_arch_left.jpg'], top=1.7,
                                    left=8.63)
    if 'resized_post_arch_upper.jpg' in processed_files:
        presentation = insert_image(presentation, 10, processed_files['resized_post_arch_upper.jpg'], top=4.05,
                                    left=2.79)
    if 'resized_post_arch_lower.jpg' in processed_files:
        presentation = insert_image(presentation, 10, processed_files['resized_post_arch_lower.jpg'], top=4.05,
                                    left=6.74)

    # Personal and Arch Images (grouped)(Pre: slide index 4, Post: slide index 11)
    # Pre personal
    if 'g_resized_pre_personal_front.jpg' in processed_files:
        presentation = insert_image(presentation, 4, processed_files['g_resized_pre_personal_front.jpg'], 0.9, 4.2)
    if 'g_resized_pre_personal_smile.jpg' in processed_files:
        presentation = insert_image(presentation, 4, processed_files['g_resized_pre_personal_smile.jpg'], 3.98, 4.2)
    if 'g_resized_pre_personal_oblique.jpg' in processed_files:
        presentation = insert_image(presentation, 4, processed_files['g_resized_pre_personal_oblique.jpg'], 7.03,
                                    4.2)
    if 'g_resized_pre_personal_profile.jpg' in processed_files:
        presentation = insert_image(presentation, 4, processed_files['g_resized_pre_personal_profile.jpg'], 10.09,
                                    4.2)

    # Post personal
    if 'g_resized_post_personal_front.jpg' in processed_files:
        presentation = insert_image(presentation, 11, processed_files['g_resized_post_personal_front.jpg'], 0.9,
                                    4.2)
    if 'g_resized_post_personal_smile.jpg' in processed_files:
        presentation = insert_image(presentation, 11, processed_files['g_resized_post_personal_smile.jpg'], 3.98,
                                    4.2)
    if 'g_resized_post_personal_oblique.jpg' in processed_files:
        presentation = insert_image(presentation, 11, processed_files['g_resized_post_personal_oblique.jpg'], 7.03,
                                    4.2)
    if 'g_resized_post_personal_profile.jpg' in processed_files:
        presentation = insert_image(presentation, 11, processed_files['g_resized_post_personal_profile.jpg'], 10.09,
                                    4.2)

    # Pre arch
    if 'g_resized_pre_arch_right.jpg' in processed_files:
        presentation = insert_image(presentation, 4, processed_files['g_resized_pre_arch_right.jpg'], top=3.46,
                                    left=2.01)
    if 'g_resized_pre_arch_front.jpg' in processed_files:
        presentation = insert_image(presentation, 4, processed_files['g_resized_pre_arch_front.jpg'], top=3.45,
                                    left=4.93)
    if 'g_resized_pre_arch_left.jpg' in processed_files:
        presentation = insert_image(presentation, 4, processed_files['g_resized_pre_arch_left.jpg'], top=3.45,
                                    left=7.84)
    if 'g_resized_pre_arch_upper.jpg' in processed_files:
        presentation = insert_image(presentation, 4, processed_files['g_resized_pre_arch_upper.jpg'], top=5.35,
                                    left=3.42)
    if 'g_resized_pre_arch_lower.jpg' in processed_files:
        presentation = insert_image(presentation, 4, processed_files['g_resized_pre_arch_lower.jpg'], top=5.35,
                                    left=6.44)

    # Post arch
    if 'g_resized_post_arch_right.jpg' in processed_files:
        presentation = insert_image(presentation, 11, processed_files['g_resized_post_arch_right.jpg'], top=3.46,
                                    left=2.01)
    if 'g_resized_post_arch_front.jpg' in processed_files:
        presentation = insert_image(presentation, 11, processed_files['g_resized_post_arch_front.jpg'], top=3.45,
                                    left=4.93)
    if 'g_resized_post_arch_left.jpg' in processed_files:
        presentation = insert_image(presentation, 11, processed_files['g_resized_post_arch_left.jpg'], top=3.45,
                                    left=7.84)
    if 'g_resized_post_arch_upper.jpg' in processed_files:
        presentation = insert_image(presentation, 11, processed_files['g_resized_post_arch_upper.jpg'], top=5.35,
                                    left=3.42)
    if 'g_resized_post_arch_lower.jpg' in processed_files:
        presentation = insert_image(presentation, 11, processed_files['g_resized_post_arch_lower.jpg'], top=5.35,
                                    left=6.44)


    # cast Images (Pre: slide index 5, Post: slide index 12)
    # pre casts
    if 'resized_pre_cast_right.jpg' in processed_files:
        presentation = insert_image(presentation, 5, processed_files['resized_pre_cast_right.jpg'], top=1.7,
                                    left=1.09)
    if 'resized_pre_cast_front.jpg' in processed_files:
        presentation = insert_image(presentation, 5, processed_files['resized_pre_cast_front.jpg'], top=1.7,
                                    left=4.87)
    if 'resized_pre_cast_left.jpg' in processed_files:
        presentation = insert_image(presentation, 5, processed_files['resized_pre_cast_left.jpg'], top=1.7,
                                    left=8.63)
    if 'resized_pre_cast_upper.jpg' in processed_files:
        presentation = insert_image(presentation, 5, processed_files['resized_pre_cast_upper.jpg'], top=4.05,
                                    left=2.8)
    if 'resized_pre_cast_lower.jpg' in processed_files:
        presentation = insert_image(presentation, 5, processed_files['resized_pre_cast_lower.jpg'], top=4.05,
                                    left=6.74)
    # post casts
    if 'resized_post_cast_right.jpg' in processed_files:
        presentation = insert_image(presentation, 12, processed_files['resized_post_cast_right.jpg'], top=1.7,
                                    left=1.09)
    if 'resized_post_cast_front.jpg' in processed_files:
        presentation = insert_image(presentation, 12, processed_files['resized_post_cast_front.jpg'], top=1.7,
                                    left=4.86)
    if 'resized_post_cast_left.jpg' in processed_files:
        presentation = insert_image(presentation, 12, processed_files['resized_post_cast_left.jpg'], top=1.7,
                                    left=8.63)
    if 'resized_post_cast_upper.jpg' in processed_files:
        presentation = insert_image(presentation, 12, processed_files['resized_post_cast_upper.jpg'], top=4.05,
                                    left=2.79)
    if 'resized_post_cast_lower.jpg' in processed_files:
        presentation = insert_image(presentation, 12, processed_files['resized_post_cast_lower.jpg'], top=4.05,
                                    left=6.74)

    # x-rays (Pre: slide index 6,7, Post: slide index 13, 14)
    # Pre-X-rays Images
    if 'resized_pre-Panoramic_X-Ray.jpg' in processed_files:
        presentation = insert_image(presentation, 6, processed_files['resized_pre-Panoramic_X-Ray.jpg'], left=1.92,
                                    top=1.33)
    if 'resized_pre-Cephalometric_X-Ray.jpg' in processed_files:
        presentation = insert_image(presentation, 7, processed_files['resized_pre-Cephalometric_X-Ray.jpg'],
                                    left=0.55, top=1.5)
    if 'resized_pre-Cephalometric_X-Ray_Tracing.jpg' in processed_files:
        presentation = insert_image(presentation, 7, processed_files['resized_pre-Cephalometric_X-Ray_Tracing.jpg'],
                                    left=6.81, top=1.5)

    # Post-X-rays Images
    if 'resized_post-Panoramic_X-Ray.jpg' in processed_files:
        presentation = insert_image(presentation, 13, processed_files['resized_post-Panoramic_X-Ray.jpg'],
                                    left=1.92, top=1.33)
    if 'resized_post-Cephalometric_X-Ray.jpg' in processed_files:
        presentation = insert_image(presentation, 14, processed_files['resized_post-Cephalometric_X-Ray.jpg'],
                                    left=0.55, top=1.5)
    if 'resized_post-Cephalometric_X-Ray_Tracing.jpg' in processed_files:
        presentation = insert_image(presentation, 14,
                                    processed_files['resized_post-Cephalometric X-Ray_Tracing.jpg'], left=6.81,
                                    top=1.5)

    # Save presentation to memory
    ppt_stream = io.BytesIO()
    presentation.save(ppt_stream)
    ppt_stream.seek(0)
    processed_files[f"{patient_name}.pptx"] = ppt_stream.getvalue()

    return processed_files

# great and show the user what is the app for
@app.route('/', methods=['GET'])
def homepage():
    return render_template('index.html')

@app.route('/doctor_info', methods=['GET', 'POST'])
def doctor_info():
    if request.method == 'POST':
        # Capture doctor info from form
        session['doc_name'] = request.form.get('doc_name')
        session['doc_email'] = request.form.get('doc_email')
        session['doc_phone'] = request.form.get('doc_phone')

        # Validate required fields
        if not all([session['doc_name'], session['doc_email'], session['doc_phone']]):
            return render_template('doctor_info.html', error="Please fill all required fields.")

        return redirect(url_for('patient_info'))

    return render_template('doctor_info.html', title='Doctor Information')

@app.route('/patient_info', methods=['GET', 'POST'])
def patient_info():
    if request.method == 'POST':
        patient_name = request.form.get('patient_name', '')
        if not patient_name:
            return redirect(url_for('patient_info'))

        # Collect uploaded files
        uploaded_files = {}
        for category in CATEGORIES.values():
            for key in category:
                file = request.files.get(key)
                if file and file.filename:
                    uploaded_files[key] = file.read()

        # Process all images
        processed_files = process_all_images(patient_name, uploaded_files)

        # Create ZIP in memory
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
            for filename, file_data in processed_files.items():
                zip_file.writestr(filename, file_data)

        zip_buffer.seek(0)

        # Send ZIP file
        return send_file(
            zip_buffer,
            mimetype='application/zip',
            as_attachment=True,
            download_name=f"{patient_name}.zip"
        )

    return render_template('patient_info.html', categories=CATEGORIES)


if __name__ == '__main__':
    app.run(debug=True)